#!/usr/bin/env python3

import os
from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title, content, note):
    """
    Adds a slide to the presentation with a title, content, and note.

    Args:
        prs (Presentation): The presentation object to which the slide will be added.
        title (str): The title of the slide.
        content (list): The content of the slide as bullet points.
        note (str): The note to be added to the slide.

    Returns:
        None
    """
    # Add a title and content slide
    slide_layout = prs.slide_layouts[1]  # 0 is for title slide, 1 for title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    # Set title
    title_placeholder.text = title
    
    # Add content as bullet points
    for i, bullet_point in enumerate(content):
        p = content_placeholder.text_frame.add_paragraph()
        p.text = bullet_point
        # p.font.size = Pt(18)
        p.level = 0
    
    # Add notes to the slide
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = note
    
    
def get_new_filename(filename):
    """
    This function takes a filename as input and returns a new filename.
    
    Parameters:
    filename (str): The original filename
    
    Returns:
    str: The new filename
    """
    base, extension = os.path.splitext(filename)
    counter = 1
    while os.path.exists(filename):
        filename = f"{base}.{counter}{extension}"
        counter += 1
    return filename

def main():
    # Define slide titles and content
    slides_content = (
        ("Who Am I?", 
        ["Platform Architect at Codence", "Skills in Python, JavaScript, C#, etc.", "Passionate about efficient coding practices"],
        "Introduce yourself as the Platform Architect at Codence. Discuss your proficiency in Python, JavaScript, and C#, and express your passion for efficient coding practices."),

        ("Target Audience and Talk Focus", 
        ["Beginner to Intermediate Level Developers", "Focus on Readability and Maintainability in Code", "Guard Clauses as a Methodology", "Ideal for Script Writers and Aspiring Software Architects"],
        "- Explain that the talk is structured to be accessible and informative for those just starting out, as well as those with some experience looking to enhance their skills.\n"
        "- Emphasize how the use of guard clauses significantly improves code readability and maintainability. You might want to give a brief example or analogy to illustrate this point.\n"
        "- Clarify that the session will delve into guard clauses, including their role within the single-pass loop methodology. Emphasize how these concepts contribute to cleaner, more understandable code.\n"
        "- Point out that the principles discussed are particularly relevant for script writers and those looking to move into software architecture, but the concepts are broadly applicable across many programming disciplines."),

        ("Understanding Guard Clauses", 
        ["Definition: Preventing Deep Nesting in Code", "Use Case: Early Exit from Functions or Loops", "Further Reading: Codementor Article, Wikipedia"],
        "- Guard clauses act as proactive measures in functions or scripts to handle special conditions or potential errors at the beginning, guarding the main logic of the code.\n"
        "- They are particularly useful for managing early exits in functions or loops, avoiding unnecessary processing by checking and handling edge cases or invalid conditions at the start.\n"
        "- Recommended resources for a deeper understanding of guard clauses include a practical Codementor article and a more theoretical Wikipedia page."),

        ("Single-Pass Loops: Emulating Try-Catch", 
        ["Understanding Single-Pass Loops", "Using 'Exit Loop If' in a Single Pass Loop", "Advantages: Clarity and Maintainability"],
        "- Describe the single-pass loop as a method to emulate try-catch patterns.\n"
        "- Explain the use of 'Exit Loop If' within this context and its advantages in terms of clarity and maintainability of code."),

        ("Error Handling: Important but Out of Scope", 
        ["The Significance of Error Handling", "Focus of This Presentation", "Resources for Learning About Error Handling"],
        "- Acknowledge the importance of error handling in programming.\n"
        "- Clarify that it's out of scope for this presentation and provide resources for attendees to explore this topic further."),

        ("Comparing IF Statements and GUARD Clauses", 
        ["Traditional IF Statements", "Guard Clauses for Early Exits", "Enhancing Readability and Maintainability"],
        "- Compare traditional IF statements with guard clauses.\n"
        "- Discuss how guard clauses can lead to early exits in functions, thus enhancing the readability and maintainability of code."),

        ("Combining 'Exit Loop If' with Single Pass Loop in FileMaker", 
        ["The Single Pass Loop Technique", "Role of 'Exit Loop If' in Nested Loops", "Practical Advantages"],
        "- Detail the combination of 'Exit Loop If' with the single-pass loop in FileMaker.\n"
        "- Explain how this approach can be used to effectively handle script logic and errors."),

        ("Nested Single Pass Loops in Iterating Loops for Try-Catch Patterns", 
        ["Embedding Single Pass Loops within Iterating Loops", "Role of 'Exit Loop If' in Nested Loops", "Practical Advantages"],
        "- Discuss embedding single-pass loops within iterating loops to create inner try-catch sequences.\n"
        "- Explain the role of 'Exit Loop If' in these nested loops and their practical advantages."),

        ("Balancing Error Handling in Scripts: How Many is Too Many?", 
        ["The Ease of Implementing Error Checks", "Finding the Balance", "Best Practices"],
        "- Discuss the ease of implementing error checks at every step of a process with guard clauses and single-pass loops.\n"
        "- Emphasize finding a balance to avoid overcomplicating scripts."),

        ("Questions & Answers", 
        ["Open Forum for Questions", "Clarifications and Further Explanations", "Sharing Practical Experiences"],
        "- Provide an opportunity for the audience to ask questions, seek clarifications, and share their own experiences related to the topics discussed."),

        ("Thank You!", 
        ["Appreciation for Attendance and Participation", "Offer of Continued Dialogue", "Contact Information: cristos.lianides-chin@codence.com"],
        "- Thank the audience for their participation.\n"
        "- Offer avenues for continued dialogue and provide your contact information for further communication.")
    )
    
    filename = "presentation.pptx"

    # Create a presentation object
    prs = Presentation()

    # Generate all slides
    for slide_title, slide_content, slide_note in slides_content:
        add_slide(prs, slide_title, slide_content, slide_note)

    # Save the presentation
    filename = get_new_filename(filename)
    prs.save(filename)

    print("Presentation created successfully!")

    # Open the presentation
    print(f"Opening {filename}...")
    os.system(f"open {filename}")  # for macOS

if __name__ == "__main__":
    main()